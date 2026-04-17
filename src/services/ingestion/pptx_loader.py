from __future__ import annotations

import io

from google.genai import types
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.logging.custom_logger import Logging

LOGGER = Logging().get_logger("pptx_loader")


def _collect_shape_text(shape) -> str:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            return text
    return ""


def pptx_to_parts(data: bytes) -> list[types.Part]:
    """Convert a PPTX file into interleaved text/image Parts for Gemini.

    For each slide: a text Part with slide number + extracted text, followed
    by image Parts for every embedded picture. Preserves both content and
    visual context without requiring LibreOffice.
    """
    prs = Presentation(io.BytesIO(data))
    parts: list[types.Part] = []
    total_images = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_chunks: list[str] = []
        images: list[tuple[bytes, str]] = []

        for shape in slide.shapes:
            text = _collect_shape_text(shape)
            if text:
                text_chunks.append(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    images.append((image.blob, image.content_type))
                except Exception as e:  # noqa: BLE001
                    LOGGER.warning(f"Slide {slide_idx}: failed to read image: {e}")

        header = f"--- Slide {slide_idx} ---"
        body = "\n".join(text_chunks) if text_chunks else "(no text)"
        parts.append(types.Part.from_text(text=f"{header}\n{body}"))

        for blob, mime in images:
            parts.append(types.Part.from_bytes(data=blob, mime_type=mime))
            total_images += 1

    LOGGER.info(
        f"Parsed PPTX: {len(prs.slides)} slides, {total_images} images, "
        f"{len(parts)} total Parts"
    )
    return parts
