import io

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from src.services.ingestion.pptx_loader import pptx_to_parts  # noqa: E402


def _build_fixture_pptx() -> bytes:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "Hello Slide 1"

    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb2.text_frame.text = "Hello Slide 2"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_to_parts_emits_one_text_part_per_slide():
    parts = pptx_to_parts(_build_fixture_pptx())
    text_parts = [p for p in parts if p.text is not None and p.text != ""]
    assert len(text_parts) == 2
    assert "Slide 1" in text_parts[0].text
    assert "Hello Slide 1" in text_parts[0].text
    assert "Slide 2" in text_parts[1].text
